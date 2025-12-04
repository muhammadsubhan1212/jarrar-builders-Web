#!/usr/bin/env python3
"""
Convert HTML presentation to PowerPoint (PPTX) format.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup

# Color definitions from the HTML
COLORS = {
    'bg_base': RGBColor(15, 23, 42),      # #0f172a
    'bg_elevated': RGBColor(2, 6, 23),    # #020617
    'bg_surface': RGBColor(30, 41, 59),   # #1e293b
    'bg_muted': RGBColor(51, 65, 85),    # #334155
    'text_primary': RGBColor(241, 245, 249),  # #f1f5f9
    'text_secondary': RGBColor(203, 213, 225),  # #cbd5e1
    'text_muted': RGBColor(148, 163, 184),  # #94a3b8
    'text_accent': RGBColor(56, 189, 248),  # #38bdf8
    'primary': RGBColor(56, 189, 248),      # #38bdf8
}

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_presentation():
    """Create a new PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)  # 960px ≈ 10 inches
    prs.slide_height = Inches(5.625)  # 540px ≈ 5.625 inches (16:9 aspect ratio)
    return prs

def add_title_slide(prs, title, subtitle, footer):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg_elevated']
    
    # Title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['text_accent']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    top = Inches(3.5)
    height = Inches(0.8)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLORS['text_secondary']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Footer
    top = Inches(4.8)
    height = Inches(0.4)
    footer_box = slide.shapes.add_textbox(left, top, width, height)
    footer_frame = footer_box.text_frame
    footer_frame.text = footer
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = COLORS['text_muted']
    footer_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, items=None, bg_color='bg_base'):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[bg_color]
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['text_accent']
    
    # Content
    if items:
        top = Inches(1.5)
        height = Inches(3.5)
        content_box = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(items):
            if i > 0:
                p = content_frame.add_paragraph()
            else:
                p = content_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text_secondary']
            p.level = 0
            p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_items=None, right_text=None, right_items=None, left_text=None, bg_color='bg_base'):
    """Add a slide with two columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[bg_color]
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['text_accent']
    
    # Determine which content goes where
    # If right_items exists, it means list is on right (like S3 slide)
    if right_items:
        # Left column: placeholder
        if left_text:
            left_col = Inches(0.5)
            top = Inches(1.5)
            width = Inches(4)
            height = Inches(3.5)
            left_box = slide.shapes.add_textbox(left_col, top, width, height)
            left_frame = left_box.text_frame
            left_frame.text = left_text
            left_para = left_frame.paragraphs[0]
            left_para.font.size = Pt(14)
            left_para.font.color.rgb = COLORS['text_muted']
            left_para.alignment = PP_ALIGN.CENTER
            
            # Add background to left box
            fill = left_box.fill
            fill.solid()
            fill.fore_color.rgb = COLORS['bg_surface']
            line = left_box.line
            line.color.rgb = COLORS['text_accent']
            line.width = Pt(2)
            line.dash_style = 2  # Dashed line
        
        # Right column: list
        right_col = Inches(5.5)
        top = Inches(1.5)
        width = Inches(4)
        height = Inches(3.5)
        right_box = slide.shapes.add_textbox(right_col, top, width, height)
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, item in enumerate(right_items):
            if i > 0:
                p = right_frame.add_paragraph()
            else:
                p = right_frame.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text_secondary']
            p.level = 0
            p.space_after = Pt(12)
    else:
        # Standard layout: list on left, placeholder on right
        # Left column (list)
        if left_items:
            left_col = Inches(0.5)
            top = Inches(1.5)
            width = Inches(4)
            height = Inches(3.5)
            left_box = slide.shapes.add_textbox(left_col, top, width, height)
            left_frame = left_box.text_frame
            left_frame.word_wrap = True
            
            for i, item in enumerate(left_items):
                if i > 0:
                    p = left_frame.add_paragraph()
                else:
                    p = left_frame.paragraphs[0]
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = COLORS['text_secondary']
                p.level = 0
                p.space_after = Pt(12)
        
        # Right column (text or placeholder)
        if right_text:
            right_col = Inches(5.5)
            top = Inches(1.5)
            width = Inches(4)
            height = Inches(3.5)
            right_box = slide.shapes.add_textbox(right_col, top, width, height)
            right_frame = right_box.text_frame
            right_frame.text = right_text
            right_para = right_frame.paragraphs[0]
            right_para.font.size = Pt(14)
            right_para.font.color.rgb = COLORS['text_muted']
            right_para.alignment = PP_ALIGN.CENTER
            
            # Add background to right box
            fill = right_box.fill
            fill.solid()
            fill.fore_color.rgb = COLORS['bg_surface']
            line = right_box.line
            line.color.rgb = COLORS['text_accent']
            line.width = Pt(2)
            line.dash_style = 2  # Dashed line
    
    return slide

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items, bg_color='bg_surface'):
    """Add a comparison slide with two cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[bg_color]
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['text_accent']
    
    # Left card
    card_left = Inches(0.5)
    card_top = Inches(1.5)
    card_width = Inches(4.25)
    card_height = Inches(3.5)
    left_card = slide.shapes.add_textbox(card_left, card_top, card_width, card_height)
    left_frame = left_card.text_frame
    left_frame.text = left_title
    left_para = left_frame.paragraphs[0]
    left_para.font.size = Pt(24)
    left_para.font.bold = True
    left_para.font.color.rgb = COLORS['text_accent']
    left_para.space_after = Pt(12)
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_secondary']
        p.level = 0
        p.space_after = Pt(8)
    
    # Add card background
    fill = left_card.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg_muted']
    line = left_card.line
    line.color.rgb = COLORS['text_primary']
    line.width = Pt(1)
    
    # Right card
    card_left = Inches(5.25)
    card_top = Inches(1.5)
    card_width = Inches(4.25)
    card_height = Inches(3.5)
    right_card = slide.shapes.add_textbox(card_left, card_top, card_width, card_height)
    right_frame = right_card.text_frame
    right_frame.text = right_title
    right_para = right_frame.paragraphs[0]
    right_para.font.size = Pt(24)
    right_para.font.bold = True
    right_para.font.color.rgb = COLORS['text_accent']
    right_para.space_after = Pt(12)
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_secondary']
        p.level = 0
        p.space_after = Pt(8)
    
    # Add card background
    fill = right_card.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg_muted']
    line = right_card.line
    line.color.rgb = COLORS['text_primary']
    line.width = Pt(1)
    
    return slide

def add_conclusion_slide(prs, title, content, footer):
    """Add a conclusion slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['bg_elevated']
    
    # Title
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['text_accent']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Content
    top = Inches(2.5)
    height = Inches(1.8)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.text = content
    content_frame.word_wrap = True
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(18)
    content_para.font.color.rgb = COLORS['text_secondary']
    content_para.alignment = PP_ALIGN.CENTER
    
    # Footer
    top = Inches(4.5)
    height = Inches(0.4)
    footer_box = slide.shapes.add_textbox(left, top, width, height)
    footer_frame = footer_box.text_frame
    footer_frame.text = footer
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = COLORS['text_muted']
    footer_para.alignment = PP_ALIGN.CENTER
    
    return slide

def parse_html_slides(html_content):
    """Parse HTML and extract slide data."""
    soup = BeautifulSoup(html_content, 'html.parser')  # Using built-in parser (no lxml needed)
    slides = []
    
    for section in soup.find_all('section', class_='slide'):
        slide_data = {'type': 'content', 'bg_color': 'bg_base'}
        
        # Get background color
        style = section.get('style', '')
        if 'background-color: #020617' in style:
            slide_data['bg_color'] = 'bg_elevated'
        elif 'background-color: #1e293b' in style:
            slide_data['bg_color'] = 'bg_surface'
        
        # Get title
        h1 = section.find('h1')
        h2 = section.find('h2')
        h3 = section.find('h3')
        if h1:
            slide_data['title'] = h1.get_text(strip=True)
        elif h2:
            slide_data['title'] = h2.get_text(strip=True)
        elif h3:
            slide_data['title'] = h3.get_text(strip=True)
        else:
            slide_data['title'] = ''
        
        # Check for center class (title/conclusion slide)
        classes = section.get('class', [])
        if 'center' in classes:
            slide_data['type'] = 'centered'
            
            # For title slide, get subtitle from text-2xl paragraph
            subtitle_p = section.find('p', class_='text-2xl')
            if subtitle_p:
                slide_data['subtitle'] = subtitle_p.get_text(strip=True)
        
        # Get list items
        ul = section.find('ul')
        if ul:
            slide_data['items'] = [li.get_text(strip=True) for li in ul.find_all('li')]
        
        # Check for two-column layout
        row = section.find('div', class_='row')
        if row:
            slide_data['type'] = 'two_column'
            fill_width_divs = row.find_all('div', class_='fill-width')
            
            # Check order: first div could be list or placeholder
            for i, div in enumerate(fill_width_divs):
                div_style = div.get('style', '')
                div_text = div.get_text()
                
                # Check if this is a placeholder (has dashed border or screenshot text)
                is_placeholder = 'dashed' in div_style or 'Screenshot' in div_text
                
                # Check if this has a list
                ul = div.find('ul')
                has_list = ul is not None
                
                if is_placeholder:
                    p = div.find('p')
                    if p:
                        if i == 0:  # First div = left side
                            slide_data['left_text'] = p.get_text(strip=True)
                        else:  # Second div = right side
                            slide_data['right_text'] = p.get_text(strip=True)
                
                if has_list:
                    items = [li.get_text(strip=True) for li in ul.find_all('li')]
                    if i == 0:  # First div = left side
                        slide_data['left_items'] = items
                    else:  # Second div = right side
                        slide_data['right_items'] = items
        
        # Check for comparison cards
        cards = section.find_all('div', class_='card-rect')
        if cards:
            slide_data['type'] = 'comparison'
            if len(cards) >= 2:
                left_card = cards[0]
                right_card = cards[1]
                left_h3 = left_card.find('h3')
                right_h3 = right_card.find('h3')
                left_ul = left_card.find('ul')
                right_ul = right_card.find('ul')
                
                slide_data['left_title'] = left_h3.get_text(strip=True) if left_h3 else ''
                slide_data['right_title'] = right_h3.get_text(strip=True) if right_h3 else ''
                slide_data['left_items'] = [li.get_text(strip=True) for li in left_ul.find_all('li')] if left_ul else []
                slide_data['right_items'] = [li.get_text(strip=True) for li in right_ul.find_all('li')] if right_ul else []
        
        # Get paragraph content for conclusion
        paragraphs = section.find_all('p', class_='text-lg')
        if paragraphs:
            # Get the first text-lg paragraph that's not in a card
            for p in paragraphs:
                if not p.find_parent('div', class_='card-rect'):
                    slide_data['content'] = p.get_text(strip=True)
                    break
        
        # Get footer
        footer_p = section.find('p', class_='text-sm')
        if footer_p:
            slide_data['footer'] = footer_p.get_text(strip=True)
        
        slides.append(slide_data)
    
    return slides

def main():
    """Main function to convert HTML to PPTX."""
    html_content = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=960, initial-scale=1.0">
  <title>Cloud Project Presentation</title>
</head>
<body class="theme-dark">
  <!-- Title Slide -->
  <section class="slide p-16 col center" style="background-color: #020617; color: #f1f5f9; display: flex; flex-direction: column; justify-content: center; align-items: center;">
    <h1 class="text-6xl mb-4" style="color: #38bdf8; margin: 0 0 1rem 0;">Cloud Project Report</h1>
    <p class="text-2xl" style="color: #cbd5e1; margin: 0;">Assessment 3: EC2, S3 & Cloud Comparison</p>
    <p class="text-sm mt-8" style="color: #94a3b8; margin-top: 2rem; margin-bottom: 0;">Presented by [Your Name] • November 2025</p>
  </section>

  <!-- Project Overview Slide -->
  <section class="slide p-16 col" style="background-color: #0f172a; color: #f1f5f9; display: flex; flex-direction: column;">
    <h2 class="mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0;">Project Overview</h2>
    <ul class="text-lg" style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
      <li>Deploy an EC2 instance and S3 bucket on AWS</li>
      <li>Demonstrate secure connection between EC2 and S3</li>
      <li>Compare AWS with Microsoft Azure cloud services</li>
      <li>Highlight core cloud computing concepts and models</li>
    </ul>
  </section>

  <!-- Cloud Computing Concepts Slide -->
  <section class="slide p-16 col" style="background-color: #0f172a; color: #f1f5f9; display: flex; flex-direction: column;">
    <h2 class="mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0;">Cloud Computing Concepts</h2>
    <ul class="text-lg" style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
      <li>On-demand self-service and scalability</li>
      <li>Infrastructure as a Service (IaaS) model</li>
      <li>Virtualisation and resource pooling principles</li>
      <li>Public cloud deployment model (AWS ownership)</li>
    </ul>
  </section>

  <!-- EC2 Instance Deployment Slide -->
  <section class="slide p-16 col" style="background-color: #0f172a; color: #f1f5f9; display: flex; flex-direction: column;">
    <h2 class="mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0;">EC2 Instance Deployment</h2>
    <div class="row gap-lg items-center" style="display: flex; flex-direction: row; gap: 1.5rem; align-items: center; flex: 1;">
      <div class="fill-width" style="flex: 1;">
        <ul class="text-lg" style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
          <li>Launched EC2 instance using Amazon Linux 2023 AMI</li>
          <li>t3.micro type, free tier, balanced CPU/memory</li>
          <li>Secure SSH access with key pair & security groups</li>
          <li>Configured storage and started instance</li>
        </ul>
      </div>
      <div class="fill-width" style="background-color: #1e293b; padding: 1.5rem; border: 2px dashed #38bdf8; text-align: center;">
        <p style="color: #94a3b8; font-size: 0.875rem;">[EC2 Instance Console Screenshot]</p>
      </div>
    </div>
  </section>

  <!-- S3 Bucket Configuration Slide -->
  <section class="slide p-16 col" style="background-color: #0f172a; color: #f1f5f9; display: flex; flex-direction: column;">
    <h2 class="mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0;">S3 Bucket Configuration</h2>
    <div class="row gap-lg items-center" style="display: flex; flex-direction: row; gap: 1.5rem; align-items: center; flex: 1;">
      <div class="fill-width" style="background-color: #1e293b; padding: 1.5rem; border: 2px dashed #38bdf8; text-align: center; flex: 1;">
        <p style="color: #94a3b8; font-size: 0.875rem; margin: 0;">[S3 Bucket Settings Screenshot]</p>
      </div>
      <div class="fill-width" style="flex: 1;">
        <ul class="text-lg" style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
          <li>Created S3 bucket in matching region (Sydney)</li>
          <li>Enabled versioning, encryption (SSE-S3)</li>
          <li>Configured tags, logging, and private access</li>
          <li>Object-based storage with high durability</li>
        </ul>
      </div>
    </div>
  </section>

  <!-- EC2 to S3 Integration Slide -->
  <section class="slide p-16 col" style="background-color: #0f172a; color: #f1f5f9; display: flex; flex-direction: column;">
    <h2 class="mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0;">EC2–S3 Integration</h2>
    <div class="row gap-lg items-center" style="display: flex; flex-direction: row; gap: 1.5rem; align-items: center; flex: 1;">
      <div class="fill-width" style="flex: 1;">
        <ul class="text-lg" style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
          <li>IAM role created (AmazonS3FullAccess) for EC2</li>
          <li>Connected EC2 to S3 with AWS CLI (secure method)</li>
          <li>Test file upload and listing confirmed integration</li>
          <li>Role-based permissions for cloud security</li>
        </ul>
      </div>
      <div class="fill-width" style="background-color: #1e293b; padding: 1.5rem; border: 2px dashed #38bdf8; text-align: center;">
        <p style="color: #94a3b8; font-size: 0.875rem;">[AWS CLI Test Output Screenshot]</p>
      </div>
    </div>
  </section>

  <!-- Cloud Provider Comparison Slide -->
  <section class="slide p-16 col" style="background-color: #0f172a; color: #f1f5f9; display: flex; flex-direction: column;">
    <h2 class="mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0;">AWS vs Azure: Comparison</h2>
    <ul class="text-lg" style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
      <li>AWS EC2 offers more instance choices than Azure VM</li>
      <li>Azure excels at Windows integration, enterprise identity</li>
      <li>S3 and Blob Storage both have strong durability, but Azure is cheaper per GB</li>
      <li>AWS IAM simpler for cloud-native, Azure AAD is advanced for enterprise</li>
    </ul>
  </section>

  <!-- Service Model & Features Slide -->
  <section class="slide p-16 col" style="background-color: #1e293b; color: #f1f5f9; display: flex; flex-direction: column;">
    <h2 class="mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0;">Service Model & Features</h2>
    <div class="row gap-lg" style="display: flex; flex-direction: row; gap: 1.5rem; flex: 1;">
      <div class="fill-width card-rect" style="background-color: #334155; color: #f1f5f9; flex: 1; padding: 2rem; box-shadow: 0 8px 16px rgba(0,0,0,0.2);">
        <h3 style="color: #38bdf8; margin: 0 0 1rem 0;">AWS Strengths</h3>
        <ul style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
          <li>Largest global footprint</li>
          <li>Developer-focused tools</li>
          <li>Linux & web application integration</li>
        </ul>
      </div>
      <div class="fill-width card-rect" style="background-color: #334155; color: #f1f5f9; flex: 1; padding: 2rem; box-shadow: 0 8px 16px rgba(0,0,0,0.2);">
        <h3 style="color: #38bdf8; margin: 0 0 1rem 0;">Azure Strengths</h3>
        <ul style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
          <li>Hybrid cloud deployment</li>
          <li>Enterprise identity (AAD, RBAC)</li>
          <li>Windows & business app integration</li>
        </ul>
      </div>
    </div>
  </section>

  <!-- Recommendations Slide -->
  <section class="slide p-16 col" style="background-color: #0f172a; color: #f1f5f9; display: flex; flex-direction: column;">
    <h2 class="mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0;">Recommendations</h2>
    <ul class="text-lg" style="color: #cbd5e1; margin: 0; padding-left: 2rem;">
      <li>For academic projects and Linux apps, AWS excels in simplicity and performance</li>
      <li>For enterprise Windows and hybrid identity needs, Azure offers superior solutions</li>
      <li>Both providers suitable for scalable, secure cloud deployments</li>
      <li>Choose based on integration demand, strategic goals</li>
    </ul>
  </section>

  <!-- Conclusion Slide -->
  <section class="slide p-16 col center" style="background-color: #020617; color: #f1f5f9; display: flex; flex-direction: column; justify-content: center; align-items: center;">
    <h2 class="text-5xl mb-6" style="color: #38bdf8; margin: 0 0 1.5rem 0; font-size: 2.25rem;">Conclusion</h2>
    <p class="text-lg" style="color: #cbd5e1; max-width: 700px; margin: 0;">
      The deployment of EC2 and S3 demonstrates cloud principles of scalability, security, and flexibility. AWS and Azure both provide robust platforms; selection should match user needs and integration context.
    </p>
    <p class="text-sm mt-8" style="color: #94a3b8; margin-top: 2rem; margin-bottom: 0;">Thank you for your attention!</p>
  </section>
</body>
</html>"""
    
    # Parse slides
    slides = parse_html_slides(html_content)
    
    # Create presentation
    prs = create_presentation()
    
    # Process each slide
    for i, slide_data in enumerate(slides):
        if slide_data['type'] == 'centered':
            if i == 0:  # Title slide
                title = slide_data.get('title', '')
                subtitle = slide_data.get('subtitle', 'Assessment 3: EC2, S3 & Cloud Comparison')
                footer = slide_data.get('footer', 'Presented by [Your Name] • November 2025')
                add_title_slide(prs, title, subtitle, footer)
            else:  # Conclusion slide
                title = slide_data.get('title', '')
                content = slide_data.get('content', '')
                footer = slide_data.get('footer', '')
                add_conclusion_slide(prs, title, content, footer)
        elif slide_data['type'] == 'two_column':
            title = slide_data.get('title', '')
            left_items = slide_data.get('left_items', [])
            right_text = slide_data.get('right_text', '')
            right_items = slide_data.get('right_items', [])
            left_text = slide_data.get('left_text', '')
            bg_color = slide_data.get('bg_color', 'bg_base')
            add_two_column_slide(prs, title, left_items, right_text, right_items, left_text, bg_color)
        elif slide_data['type'] == 'comparison':
            title = slide_data.get('title', '')
            left_title = slide_data.get('left_title', '')
            left_items = slide_data.get('left_items', [])
            right_title = slide_data.get('right_title', '')
            right_items = slide_data.get('right_items', [])
            bg_color = slide_data.get('bg_color', 'bg_surface')
            add_comparison_slide(prs, title, left_title, left_items, right_title, right_items, bg_color)
        else:  # Regular content slide
            title = slide_data.get('title', '')
            items = slide_data.get('items', [])
            bg_color = slide_data.get('bg_color', 'bg_base')
            add_content_slide(prs, title, items, bg_color)
    
    # Save presentation
    output_file = 'Cloud_Project_Presentation.pptx'
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == '__main__':
    main()

